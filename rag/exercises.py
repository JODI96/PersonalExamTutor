"""
Konvertiert echte Kursübungen aus .adoc / .pptx / .pdf Dateien in Typst.
Lösungen werden mit Claude detaillierter gemacht.
"""

from __future__ import annotations

import os
import re
from pathlib import Path

import anthropic
from dotenv import load_dotenv

from rag.config import PDFS_DIR, DOCS_WEEKS, EXERCISES_DIR, CLAUDE_MODEL
from rag.summarize import _build_pdf, _ensure_summary_includes

load_dotenv()

_SUPPORTED_EXT = {".adoc", ".pdf", ".pptx", ".txt", ".md"}


def _files_in(folder: Path) -> list[Path]:
    if not folder.exists():
        return []
    return sorted(f for f in folder.iterdir()
                  if f.is_file() and f.suffix.lower() in _SUPPORTED_EXT)


# ── Text-Extraktion ───────────────────────────────────────────────────────────

def _read_adoc(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return f"[PPTX-Inhalt konnte nicht gelesen werden: python-pptx nicht installiert]"
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides.append(f"--- Folie {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _read_pdf(path: Path) -> str:
    import fitz
    doc = fitz.open(str(path))
    pages = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(pages)


def _read_any(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".adoc":
        return _read_adoc(path)
    if ext == ".pptx":
        return _read_pptx(path)
    if ext == ".pdf":
        return _read_pdf(path)
    return path.read_text(encoding="utf-8", errors="replace")


# ── Claude-Prompts ────────────────────────────────────────────────────────────

_EXERCISES_TO_TYPST_SYSTEM = """\
Du bist ein Typst-Formatierer. Du erhältst echte Kursübungen im AsciiDoc- oder \
Plaintext-Format und wandelst sie in sauberes Typst um.

Struktur:
  == Übungen

  === Aufgabe 1 — [Titel falls vorhanden]
  [Aufgabentext]

  **(a)** [Teilaufgabe]
  **(b)** [Teilaufgabe]
  ...

Regeln:
- Übernimm den Aufgabentext EXAKT — erfinde keine neuen Aufgaben
- Wandle AsciiDoc-Syntax in Typst: `*fett*` → `*fett*`, ` `code` ` → `\`code\``
- Mathematik: $...$ für inline, $...$ für Display (Typst-Syntax, kein LaTeX)
- Typst-Mathe: frac(a,b), bold(x), sum_(i=0)^n, mat(1,2;3,4) usw.
- Kein #import, kein #show — beginne direkt mit == Übungen
- Nummeriere Aufgaben falls noch nicht nummeriert
"""

_SOLUTIONS_ENRICH_SYSTEM = """\
Du bist ein Prüfungstutor. Du erhältst:
- Die Vorlesungsunterlagen der Woche (Folien/Notizen)
- Die Übungsaufgaben
- Vorhandene Musterlösungen (falls vorhanden)

Deine Aufgabe: Erstelle vollständig detaillierte Schritt-für-Schritt Lösungen.

Wichtig: Nutze die TERMINOLOGIE und NOTATION aus den Vorlesungsunterlagen — \
erkläre so, wie es der Professor in den Folien erklärt hat. \
Verweise auf Konzepte aus der Vorlesung wo passend.

Schritte:
1. Übernimm die vorhandene Lösung als Grundlage (falls vorhanden)
2. Füge ALLE fehlenden Zwischenschritte ein
3. Erkläre bei jedem Schritt in einem Satz WARUM dieser Schritt gemacht wird
4. Zeige vollständige Rechnungen (keine "..." oder "analog")
5. Hebe das Endergebnis jeder Teilaufgabe klar hervor

Ausgabeformat: Typst

Struktur:
  == Lösungen

  === Aufgabe 1 — [Titel]

  ==== (a)
  *Schritt 1 —* [Erklärung]
  $ [Formel] $
  *Schritt 2 —* ...
  *Ergebnis:* $ [Endergebnis] $

Regeln:
- Typst-Mathe korrekt: frac(a,b) nicht frac(a}{b), Klammern rund nicht geschwungen
- Kein #import — beginne direkt mit == Lösungen
- Antworte auf Deutsch
"""


# ── Typst-Pfad (nur als Referenz-Datei, nicht ins PDF eingebunden) ────────────

def _exercises_typ_path(week_nr: int) -> Path:
    return DOCS_WEEKS / f"week_{week_nr:02d}" / f"exercises_{week_nr:02d}.typ"


# ── Haupt-Funktion ────────────────────────────────────────────────────────────

def convert_exercises(week_nr: int, progress=None) -> dict:
    """
    Liest echte Übungen aus readme.adoc + losung.adoc (oder PPTX/PDF),
    konvertiert zu Typst, reichert Lösungen an.

    Returns:
        {"exercises_path": str, "solutions_path": str, "sources": list[str]}
    """
    week_path = PDFS_DIR / f"week_{week_nr:02d}"
    if not week_path.exists():
        raise FileNotFoundError(f"Ordner nicht gefunden: {week_path}")

    # ── Dateien aus fixen Unterordnern lesen ──────────────────────────────────
    if progress:
        progress("Lese Dateien aus uebungen/ und loesungen/...")

    ex_files      = _files_in(week_path / "uebungen")
    sol_files     = _files_in(week_path / "loesungen")
    lecture_files = _files_in(week_path / "vorlesung")

    if not ex_files:
        raise FileNotFoundError(
            f"Keine Dateien in {week_path / 'uebungen'} gefunden.\n"
            f"Lege deine Übungs-PDFs / .adoc Dateien in diesen Ordner."
        )

    sources = [f.name for f in ex_files] + [f.name for f in sol_files]

    if progress:
        progress(f"Vorlesung: {', '.join(f.name for f in lecture_files) if lecture_files else 'keine'}")
        progress(f"Übungen:   {', '.join(f.name for f in ex_files)}")
        progress(f"Lösungen:  {', '.join(f.name for f in sol_files) if sol_files else 'keine'}")

    lecture_raw   = "\n\n---\n\n".join(_read_any(f) for f in lecture_files)
    exercises_raw = "\n\n---\n\n".join(_read_any(f) for f in ex_files)
    solutions_raw = "\n\n---\n\n".join(_read_any(f) for f in sol_files)

    client = anthropic.Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])

    # ── Übungen → Typst ───────────────────────────────────────────────────────
    if progress:
        progress("Konvertiere Übungen zu Typst...")

    ex_msg = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=4096,
        system=_EXERCISES_TO_TYPST_SYSTEM,
        messages=[{
            "role": "user",
            "content": (
                f"Woche {week_nr} Übungen:\n\n{exercises_raw}\n\n"
                f"Konvertiere diese Aufgaben in sauberes Typst."
            ),
        }],
    )
    exercises_typst = ex_msg.content[0].text.strip()

    # ── Lösungen anreichern ───────────────────────────────────────────────────
    lecture_section = (
        f"Vorlesungsunterlagen (Woche {week_nr}):\n\n{lecture_raw}\n\n"
        if lecture_raw else ""
    )

    if solutions_raw:
        if progress:
            progress("Reichere Lösungen mit detaillierten Schritten an (inkl. Vorlesungskontext)...")
        user_content = (
            f"{lecture_section}"
            f"Aufgaben (Woche {week_nr}):\n\n{exercises_raw}\n\n"
            f"Vorhandene Lösung:\n\n{solutions_raw}\n\n"
            f"Erstelle eine vollständig detaillierte Schritt-für-Schritt Lösung "
            f"in der Terminologie der Vorlesung."
        )
    else:
        if progress:
            progress("Keine Lösungsdatei gefunden — generiere Lösungen aus Aufgaben (inkl. Vorlesungskontext)...")
        user_content = (
            f"{lecture_section}"
            f"Aufgaben (Woche {week_nr}):\n\n{exercises_raw}\n\n"
            f"Erstelle vollständige detaillierte Schritt-für-Schritt Lösungen "
            f"in der Terminologie der Vorlesung."
        )

    sol_msg = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=8192,
        system=_SOLUTIONS_ENRICH_SYSTEM,
        messages=[{"role": "user", "content": user_content}],
    )
    solutions_typst = sol_msg.content[0].text.strip()

    # ── In Dateien schreiben ──────────────────────────────────────────────────
    from datetime import date
    header = (
        f"// Echte Kursübungen — Woche {week_nr} (konvertiert am {date.today()})\n"
        f"// Quelle: {', '.join(sources)}\n\n"
        f'#import "../../template.typ": tip-box, warn-box, formula-box, cheat-box\n\n'
    )

    ex_typ_path = _exercises_typ_path(week_nr)
    ex_typ_path.parent.mkdir(parents=True, exist_ok=True)
    ex_typ_path.write_text(header + exercises_typst + "\n\n" + solutions_typst, encoding="utf-8")

    # Auch als Markdown speichern (für Git-Lesbarkeit)
    ex_dir = EXERCISES_DIR / f"week_{week_nr:02d}"
    ex_dir.mkdir(parents=True, exist_ok=True)
    (ex_dir / "problems.md").write_text(
        f"# Week {week_nr} — Exercises\n\n"
        f"> Quelle: {', '.join(sources)}\n\n---\n\n{exercises_raw}",
        encoding="utf-8",
    )
    (ex_dir / "solutions.md").write_text(
        f"# Week {week_nr} — Solutions (detailliert)\n\n---\n\n{solutions_typst}",
        encoding="utf-8",
    )

    if progress:
        progress("Baue PDF neu...")
    _build_pdf()

    return {
        "exercises_path": str(ex_typ_path),
        "solutions_path": str(ex_dir / "solutions.md"),
        "sources": sources,
    }
